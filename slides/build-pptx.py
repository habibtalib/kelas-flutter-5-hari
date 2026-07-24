#!/usr/bin/env python
"""Build an editable .pptx from the reveal.js deck (flutter-training.html).

Maps the deck's design-system vocabulary (see _build/_SPEC.md) to PowerPoint
shapes, keeping the KPT navy/gold theme. Speaker notes are carried over into
the PowerPoint notes pane.

Usage:
    python build-pptx.py            # writes flutter-training.pptx next to the html
"""
from pathlib import Path

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
SRC = HERE / "flutter-training.html"
OUT = HERE / "flutter-training.pptx"

# --- KPT theme (mirrors _build/00-head.html :root) ---------------------------
NAVY = RGBColor(0x1A, 0x2B, 0x5C)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
INK = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x56, 0x63, 0x71)
LINE = RGBColor(0xE2, 0xE5, 0xEE)
PAPER = RGBColor(0xF5, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODEBG = RGBColor(0x1E, 0x22, 0x2C)
CODEFG = RGBColor(0xE6, 0xE8, 0xEE)
GREEN = RGBColor(0x16, 0x79, 0x4A)
RED = RGBColor(0xB4, 0x23, 0x18)
AMBER = RGBColor(0xB4, 0x53, 0x09)

SANS = "Segoe UI"
MONO = "Consolas"

SW, SH = 13.333, 7.5          # 16:9
ML, MR, MT = 0.70, 0.70, 0.55
CW = SW - ML - MR
BOTTOM = SH - 0.45            # content must stop above this

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# --- helpers -----------------------------------------------------------------
def text_of(el):
    """Visible text of an element, with <br> as newlines, whitespace collapsed."""
    if el is None:
        return ""
    clone = BeautifulSoup(str(el), "html.parser")
    for br in clone.find_all("br"):
        br.replace_with("\n")
    return " ".join(clone.get_text().split())


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def textbox(slide, l, t, w, h, lines, size, color, *, bold=False,
            align=PP_ALIGN.LEFT, font=SANS, anchor=MSO_ANCHOR.TOP, spacing=1.05):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(2)
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def rrect(slide, l, t, w, h, fill, line=None, radius=0.10):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    sp.text_frame.word_wrap = True
    return sp


def est_lines(s, chars_per_line):
    """Rough wrapped-line count for height estimation."""
    n = 0
    for para in (s or "").split("\n"):
        n += max(1, (len(para) // max(1, chars_per_line)) + 1)
    return n


# --- block renderers ---------------------------------------------------------
def render_bullets(slide, y, items, size=15):
    h = 0.0
    for it in items:
        t = text_of(it)
        if not t:
            continue
        lines = est_lines(t, int(CW * 11))
        bh = 0.245 * lines
        if y + bh > BOTTOM:
            break
        textbox(slide, ML + 0.18, y, CW - 0.18, bh, f"•  {t}", size, MUTED)
        y += bh
        h += bh
    return y + 0.08


def render_code(slide, y, code_text, size=11):
    lines = code_text.split("\n")
    # keep slides readable: cap very long blocks
    if len(lines) > 18:
        lines = lines[:18] + ["…"]
    bh = min(0.205 * len(lines) + 0.22, BOTTOM - y)
    if bh < 0.4:
        return y
    rrect(slide, ML, y, CW, bh, CODEBG, radius=0.06)
    tb = slide.shapes.add_textbox(Inches(ML + 0.14), Inches(y + 0.10),
                                  Inches(CW - 0.28), Inches(bh - 0.20))
    tf = tb.text_frame
    tf.word_wrap = False
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.name = MONO
        r.font.color.rgb = CODEFG
    return y + bh + 0.12


def render_table(slide, y, table, size=11):
    rows = table.find_all("tr")
    if not rows:
        return y
    grid = []
    for tr in rows:
        cells = tr.find_all(["th", "td"])
        if cells:
            grid.append([text_of(c) for c in cells])
    if not grid:
        return y
    ncols = max(len(r) for r in grid)
    grid = [r + [""] * (ncols - len(r)) for r in grid]
    nrows = len(grid)
    rh = 0.30
    th = min(rh * nrows, BOTTOM - y)
    if th < 0.4:
        return y
    shape = slide.shapes.add_table(nrows, ncols, Inches(ML), Inches(y),
                                   Inches(CW), Inches(th))
    tbl = shape.table
    for ri, row in enumerate(grid):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            cell.margin_left = cell.margin_right = Pt(5)
            cell.margin_top = cell.margin_bottom = Pt(2)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size)
                    r.font.name = SANS
                    r.font.bold = (ri == 0)
                    r.font.color.rgb = WHITE if ri == 0 else MUTED
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if ri == 0 else WHITE
    return y + th + 0.14


def render_cards(slide, y, cards, three=False):
    n = len(cards)
    if not n:
        return y
    per_row = 3 if three else 2
    gap = 0.20
    cw = (CW - gap * (per_row - 1)) / per_row
    rows = (n + per_row - 1) // per_row
    ch = 1.02
    total = rows * (ch + gap)
    if y + total > BOTTOM:
        ch = max(0.6, (BOTTOM - y) / rows - gap)
        total = rows * (ch + gap)
    for i, card in enumerate(cards):
        r, c = divmod(i, per_row)
        cx = ML + c * (cw + gap)
        cy = y + r * (ch + gap)
        if cy + ch > BOTTOM:
            break
        classes = card.get("class", [])
        navy_card = "navy" in classes
        accent = "accent" in classes
        fill = NAVY if navy_card else WHITE
        border = None if navy_card else (GOLD if accent else LINE)
        rrect(slide, cx, cy, cw, ch, fill, border, radius=0.10)
        head = card.find("h4")
        body = card.find("p")
        ty = cy + 0.10
        if head:
            textbox(slide, cx + 0.14, ty, cw - 0.28, 0.30, text_of(head), 13,
                    WHITE if navy_card else NAVY, bold=True)
            ty += 0.32
        if body:
            textbox(slide, cx + 0.14, ty, cw - 0.28, ch - (ty - cy) - 0.08,
                    text_of(body), 11, WHITE if navy_card else MUTED)
    return y + total + 0.06


def render_note(slide, y, note):
    classes = note.get("class", [])
    accent = GOLD if "gold" in classes else RED if "warn" in classes else \
        GREEN if "tip" in classes else NAVY
    t = text_of(note)
    if not t:
        return y
    nh = 0.26 * est_lines(t, int(CW * 12)) + 0.20
    if y + nh > BOTTOM:
        nh = BOTTOM - y
    if nh < 0.35:
        return y
    rrect(slide, ML, y, CW, nh, PAPER, radius=0.06)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(y),
                                 Inches(0.06), Inches(nh))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False
    textbox(slide, ML + 0.20, y + 0.08, CW - 0.34, nh - 0.16, t, 12, INK)
    return y + nh + 0.12


def render_pills(slide, y, pills):
    """Render a row of pills as one wrapped line (keeps it simple & editable)."""
    labels = [text_of(p) for p in pills if text_of(p)]
    if not labels:
        return y
    t = "   ".join(f"[ {l} ]" for l in labels)
    h = 0.26 * est_lines(t, int(CW * 11)) + 0.06
    if y + h > BOTTOM:
        return y
    textbox(slide, ML, y, CW, h, t, 12, NAVY, bold=True)
    return y + h + 0.06


def render_steps(slide, y, steps, size=14):
    for i, li in enumerate(steps, 1):
        t = text_of(li)
        if not t:
            continue
        lines = est_lines(t, int(CW * 11))
        bh = 0.26 * lines
        if y + bh > BOTTOM:
            break
        textbox(slide, ML + 0.10, y, CW - 0.10, bh, f"{i}.  {t}", size, MUTED)
        y += bh
    return y + 0.08


def render_blocks(slide, y, container, depth=0):
    """Walk direct children of a container and render known block types."""
    for el in container.children:
        if isinstance(el, NavigableString) or not isinstance(el, Tag):
            continue
        if y >= BOTTOM - 0.25:
            break
        cls = el.get("class", [])
        name = el.name

        if name == "aside":                      # speaker notes handled elsewhere
            continue
        if "slide-tag" in cls:
            continue
        if name in ("h1", "h2", "h3", "h4"):
            size = 26 if name in ("h1", "h2") else 18
            t = text_of(el)
            h = 0.40 if name in ("h1", "h2") else 0.34
            textbox(slide, ML, y, CW, h, t, size, NAVY, bold=True)
            y += h + 0.08
        elif name == "p":
            t = text_of(el)
            if not t:
                continue
            lead = "lead" in cls
            small = "small" in cls
            size = 17 if lead else (11 if small else 14)
            color = INK if lead else MUTED
            h = 0.27 * est_lines(t, int(CW * (10 if lead else 12)))
            textbox(slide, ML, y, CW, h, t, size, color, bold=lead)
            y += h + 0.08
        elif name in ("ul", "ol"):
            if "steps" in cls:
                y = render_steps(slide, y, el.find_all("li", recursive=False))
            else:
                y = render_bullets(slide, y, el.find_all("li", recursive=False))
        elif name == "pre":
            code = el.find("code")
            y = render_code(slide, y, code.get_text() if code else el.get_text())
        elif name == "table":
            y = render_table(slide, y, el)
        elif name == "div" and "cards" in cls:
            y = render_cards(slide, y, el.find_all("div", class_="card",
                                                   recursive=False),
                             three="c3" in cls)
        elif name == "div" and "note" in cls:
            y = render_note(slide, y, el)
        elif name == "div" and ("two-col" in cls or "three-col" in cls):
            # flatten columns vertically — keeps text editable and readable
            for col in el.find_all("div", recursive=False):
                y = render_blocks(slide, y, col, depth + 1)
        elif name == "div":
            y = render_blocks(slide, y, el, depth + 1)
        elif name == "span" and "pill" in cls:
            continue  # handled by the pill sweep below
    return y


def add_notes(slide, section):
    aside = section.find("aside", class_="notes")
    if aside:
        t = text_of(aside)
        if t:
            slide.notes_slide.notes_text_frame.text = t


# --- main --------------------------------------------------------------------
def main():
    soup = BeautifulSoup(SRC.read_text(encoding="utf-8"), "html.parser")
    slides_div = soup.find("div", class_="slides")
    sections = slides_div.find_all("section", recursive=False)
    print(f"parsing {len(sections)} sections from {SRC.name}")

    for idx, sec in enumerate(sections, 1):
        cls = sec.get("class", [])
        slide = prs.slides.add_slide(BLANK)

        if "title-slide" in cls or "chapter" in cls:
            set_bg(slide, NAVY)
            y = 1.45 if "title-slide" in cls else 1.70
            badge = sec.find(class_="badge")
            kicker = sec.find(class_="kicker")
            if badge:
                textbox(slide, ML, 0.85, CW, 0.75, text_of(badge), 40, GOLD)
                y = 1.85
            if kicker:
                textbox(slide, ML, y - 0.55, CW, 0.34, text_of(kicker).upper(),
                        13, GOLD, bold=True)
            head = sec.find(["h1", "h2"])
            if head:
                t = text_of(head)
                size = 40 if "title-slide" in cls else 34
                h = 0.62 * est_lines(t, 34)
                textbox(slide, ML, y, CW, h, t, size, WHITE, bold=True)
                y += h + 0.12
            sub = sec.find(class_="sub")
            if sub:
                textbox(slide, ML, y, CW, 0.42, text_of(sub), 18, GOLD, bold=True)
                y += 0.50
            agenda = sec.find("ul", class_="agenda")
            if agenda:
                for li in agenda.find_all("li", recursive=False):
                    textbox(slide, ML + 0.10, y, CW, 0.30,
                            f"•  {text_of(li)}", 13, RGBColor(0xD3, 0xD9, 0xE6))
                    y += 0.30
            meta = sec.find(class_="meta")
            if meta:
                textbox(slide, ML, SH - 1.15, CW, 0.40, text_of(meta), 12,
                        RGBColor(0xB9, 0xC2, 0xD4))
        else:
            set_bg(slide, WHITE)
            y = MT
            head = sec.find(["h2", "h3"])
            if head:
                t = text_of(head)
                h = 0.46 * est_lines(t, 46)
                textbox(slide, ML, y, CW, h, t, 24, NAVY, bold=True)
                y += h + 0.14
                head.extract()          # don't render twice
            y = render_blocks(slide, y, sec)
            # any pills not inside a rendered block
            loose = [p for p in sec.find_all("span", class_="pill")]
            if loose and y < BOTTOM - 0.3:
                render_pills(slide, y, loose)

        add_notes(slide, sec)

        # slide number (skip on title)
        if "title-slide" not in cls:
            textbox(slide, SW - 1.15, SH - 0.42, 0.75, 0.26, str(idx), 10,
                    RGBColor(0x9A, 0xA3, 0xB2), align=PP_ALIGN.RIGHT)

    prs.save(OUT)
    print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)"
          if False else f"wrote {OUT}")
    print(f"slides: {len(sections)}")


if __name__ == "__main__":
    main()
